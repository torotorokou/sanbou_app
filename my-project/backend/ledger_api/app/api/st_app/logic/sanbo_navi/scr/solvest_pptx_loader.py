# logic/sanbo_navi/scr/solvest_pptx_loader.py

import os
from typing import List, Dict, Any
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont
from langchain.schema import Document

def load_pptx_as_documents(pptx_path: str, source_id: str = "SOLVEST_PPTX") -> List[Document]:
    """
    PowerPointファイル（.pptx）をスライド単位でLangChain Documentに構造化
    """
    if not os.path.exists(pptx_path):
        raise FileNotFoundError(f"❌ 指定されたpptxが存在しません: {pptx_path}")

    prs = Presentation(pptx_path)
    documents: List[Document] = []

    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text.strip())

        content = "\n".join([t for t in texts if t])
        if not content.strip():
            continue

        metadata: Dict[str, Any] = {
            "source": source_id,
            "page": str(i + 1),
            "title": texts[0] if texts else f"Slide {i + 1}",
            "category": ["事業計画(PPTX)"],
            "tag": ["pptx", "slide"]
        }
        doc = Document(page_content=content, metadata=metadata)
        documents.append(doc)

    return documents

def render_pptx_slide_textonly(pptx_path: str, slide_number: int, width: int = 960, height: int = 540) -> Image.Image:
    """
    PPTXスライドを簡易的に画像化（テキストのみ描画）
    """
    prs = Presentation(pptx_path)
    slide_index = max(0, min(slide_number - 1, len(prs.slides) - 1))
    slide = prs.slides[slide_index]

    img = Image.new("RGB", (width, height), color="white")
    draw = ImageDraw.Draw(img)
    y = 40
    title_font = ImageFont.load_default()
    content_font = ImageFont.load_default()

    if slide.shapes.title:
        draw.text((30, y), slide.shapes.title.text.strip(), fill="black", font=title_font)
        y += 30

    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text and shape != slide.shapes.title:
            for line in shape.text.split("\n"):
                draw.text((50, y), line.strip(), fill="black", font=content_font)
                y += 20
                if y > height - 40:
                    break

    return img